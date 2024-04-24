from copy import deepcopy
import pptx
import os
import pandas as pd
from pptx.chart.data import ChartData
import requests
import wget
from responseReading import mainBrand,noOfComp
from downloadFileFromDrive import downloadFiles
from sendDeck import send
from dotenv import load_dotenv 

presentation = pptx.Presentation(os.path.join("Pixability", "input.pptx"))
load_dotenv()
pptLogo = os.path.join("logo.png")
downloadFiles()
excelPathChannel = os.path.join("Excel","channel.xlsx")
excelPathVideo = os.path.join("Excel","video.xlsx")
os.environ['DB_ADDR'] = 'rethinkdb'
apiKey = os.getenv('API_KEY')
print(apiKey)
# apiKey = "AIzaSyAevsj4BhzT6fQybkfIT4t46qkLVTIzv7Q"


def changeTextType(fontType):
    for slide in presentation.slides:
            for shape in slide.shapes:
                # Changing para font type
                if hasattr(shape, "text"):
                    for para in shape.text_frame.paragraphs:
                        if(para.runs):
                            for run in para.runs:
                                run.font.name=fontType
                # Changing table font type
                if(shape.has_table):
                    table=shape.table
                    for r in table.rows:
                        for cell in r.cells:
                            cpara = cell.text_frame.paragraphs[0]
                            if(cpara.runs):
                                for crun in cpara.runs:
                                    crun.font.name=fontType
                if(shape.has_chart):
                    chart=shape.chart
                    chart.legend.font.name = fontType
                    chart.chart_title.text_frame.paragraphs[0].runs[0].font.name = fontType
                    # axis is remaining
    

def addComa(slideNumber,rs,re,cs,ce):
    slide = presentation.slides[slideNumber]
    for shape in slide.shapes:
        if(shape.has_table):
            table=shape.table
            for r in range(rs,re+1):
                for c in range(cs,ce+1):
                    cell_tf = (table.cell(r,c)).text_frame
                    cell_value = cell_tf.text
                    temp=""
                    if(len(cell_value)<=3):
                        for x in range(len(cell_value)-1,-1,-1):
                            temp+=cell_value[x]
                    if(len(cell_value)>3):
                        for x in range(len(cell_value)-1,len(cell_value)-4,-1):
                            temp+=cell_value[x]
                        temp+=","
                        if(len(cell_value)<=5):
                            temp+=cell_value[x-1]
                        if(len(cell_value)==5):
                            temp+=cell_value[x-2]
                    if(len(cell_value)>5):
                        count=0
                        for y in range(len(cell_value)-4,-1,-1):
                            if(count==2):
                                temp+=","
                                count=0
                            temp+=cell_value[y]
                            count+=1
                    temp=temp[::-1]
                    run=cell_tf.paragraphs[0].runs[0]
                    run.text = run.text.replace(cell_value, temp)
            
    for slide in presentation.slides:
        for shape in slide.shapes:
            if(shape.has_table):
                    table=shape.table
                    for r in table.rows:
                        for cell in r.cells:
                            celltext = cell.text_frame.paragraphs[0].text
                            celltext=celltext.split(" ")[0]
                            
def addPptLogo(pptLogo,slideNumber,imgNumber):
    slide = presentation.slides[slideNumber]
    images=[]
    for shape in slide.shapes:
        if(shape.shape_type == 13):
            images.append(shape)
    img_shape = images[imgNumber]
    slide.shapes.add_picture(pptLogo, img_shape.left, img_shape.top,img_shape.width,img_shape.height)
    slide.shapes._spTree.remove(img_shape._element)

def makePivot(excelPathChannel,mainBrand):
    df = pd.read_excel(excelPathChannel,"channel")
    print(df.head(5))
    df = df.sort_values("Date",ascending=False).reset_index(drop=True)
    lastDate=df["Date"][0]
    pt = pd.pivot_table(df[df["Date"]==lastDate],values = ["SUBSCRIBERS","VIDEO_COUNT","ORGANIC_VIEWS","PAID_VIEWS","TRUEVIEW_SPEND_ESTIMATE"],
                        index=["CHANNEL_NAME"],
                        aggfunc={"SUBSCRIBERS":"sum","VIDEO_COUNT":"sum","ORGANIC_VIEWS":"sum","PAID_VIEWS":"sum","TRUEVIEW_SPEND_ESTIMATE":"sum"},  
                        )
    pt["Total View"]=pt["ORGANIC_VIEWS"]+pt["PAID_VIEWS"]
    pt = pt.iloc[:,[2,4,0,1,5,3]]
    firstRow = pt[pt.index==mainBrand]
    remainingRow = pt[pt.index!=mainBrand].sort_values("SUBSCRIBERS",ascending=False)
    newTable = pd.concat([firstRow,remainingRow])
    return newTable

def makePivotVideo(excelPathVideo,mainBrand):
    df = pd.read_excel(excelPathVideo,"video")
    pt = pd.pivot_table(df,values = ["VIEWS"],
                        index=["IAB_TIER1"],
                        aggfunc={"VIEWS":"sum"},
                        columns=["CHANNEL_NAME"]  
                        )
    return pt

def fetchData5(excelPathChannel,mainBrand,slideNo):
    newTable = makePivot(excelPathChannel,mainBrand)
    arr = newTable.to_numpy().astype(int)
    slide = presentation.slides[slideNo-1]
    cnt=1
    for shape in slide.shapes:
        if(shape.has_table):
            table=shape.table
            for brands in newTable.index:
                run = (table.cell(cnt,0)).text_frame.paragraphs[0].runs[0]
                run.text = run.text.replace(run.text,brands)
                cnt+=1
            for r in range(1,(2+noOfComp)):
                for c in range(1,6):
                    run = (table.cell(r,c)).text_frame.paragraphs[0].runs[0]
                    run.text = run.text.replace(run.text,arr[r-1][c-1].astype(str))
            temp=1
            reqHeight=0
            for row in table.rows:
                if(temp==1):
                    firstRowHeight=row.height
                reqHeight+=row.height
                temp+=1
            tbl = table._tbl
            tobeRemoved=4-noOfComp
            if(tobeRemoved!=0):
                for x in range(1,tobeRemoved+1):
                    tr = table.rows[6-x]._tr
                    tbl.remove(tr)
            temp=1
            for row in table.rows:
                if(temp==1):
                    row.height=firstRowHeight
                else:
                    row.height=(reqHeight-firstRowHeight)//(1+noOfComp)
                temp+=1
    return newTable.index

def fetchData6(excelPathChannel,mainBrand,slideNo,sortedBrands):
    newTable = makePivot(excelPathChannel,mainBrand)
    newTable['VIDEO_COUNT_PER'] = (newTable['VIDEO_COUNT'] / newTable['VIDEO_COUNT'].sum()).round(4)
    newTable['SUBSCRIBERS_PER'] = (newTable['SUBSCRIBERS'] / newTable['SUBSCRIBERS'].sum()).round(4)
    slide = presentation.slides[slideNo-1]
    for shape in slide.shapes:
        if(shape.has_chart):
            chart=shape.chart
            title = chart.chart_title.text_frame.text
            if(title.__contains__("Video")==True):
                chart_data = ChartData()
                chart_data.categories = newTable.index
                chart_data.add_series('Video Cout',newTable['VIDEO_COUNT_PER'])
                chart.replace_data(chart_data)
            if(title.__contains__("Subscriber")==True):
                chart_data = ChartData()
                chart_data.categories = newTable.index
                chart_data.add_series('Subcriber Cout',newTable['SUBSCRIBERS_PER'])
                chart.replace_data(chart_data)

def fetchData7(excelPathChannel,mainBrand,slideNo,sortedBrands):
    newTable = makePivot(excelPathChannel,mainBrand)
    newTable['TRUEVIEW_SPEND_ESTIMATE_PER'] = (newTable['TRUEVIEW_SPEND_ESTIMATE'] / newTable['TRUEVIEW_SPEND_ESTIMATE'].sum()).round(4)
    newTable['Total View_PER'] = (newTable['Total View'] / newTable['Total View'].sum()).round(4)
    slide = presentation.slides[slideNo-1]
    for shape in slide.shapes:
        if(shape.has_chart):
            chart=shape.chart
            chart_data = ChartData()
            chart_data.categories = newTable.index
            chart_data.add_series('Views Share',newTable['Total View_PER'])
            chart_data.add_series('Spend Share',newTable['TRUEVIEW_SPEND_ESTIMATE_PER'])
            chart.replace_data(chart_data)

def fetchData10(excelPathChannel,mainBrand,slideNo,sortedBrands):
    newTable = makePivot(excelPathChannel,mainBrand)
    newTable['Organic_Per'] = (newTable['ORGANIC_VIEWS'] / (newTable['ORGANIC_VIEWS']+newTable['PAID_VIEWS'])).round(4)
    newTable['Paid_Per'] = (newTable['PAID_VIEWS'] / (newTable['ORGANIC_VIEWS']+newTable['PAID_VIEWS'])).round(4)
    slide = presentation.slides[slideNo-1]
    for shape in slide.shapes:
        if(shape.has_chart):
            chart=shape.chart
            chart_data = ChartData()
            chart_data.categories = sortedBrands
            chart_data.add_series('Organic Views',newTable['Organic_Per'])
            chart_data.add_series('Paid Views',newTable['Paid_Per'])
            chart.replace_data(chart_data)
    
def fetchData11(excelPathVideo,mainBrand,slideNo,sortedBrands):
    newTable = makePivotVideo(excelPathVideo,mainBrand)
    newTable.columns=newTable.columns.droplevel(level=0)
    arr = []
    for brand in sortedBrands:
        newTable[brand+"_PER"] = (newTable[brand] / newTable[brand].sum()*100).round(2)
        df = newTable.sort_values(brand+"_PER",ascending=False)[brand+"_PER"].dropna().reset_index().iloc[:2].to_numpy()
        arr.append(df[0])
        arr.append(df[1])
    slide = presentation.slides[slideNo-1]
    cnt=1
    for shape in slide.shapes:
        if(shape.has_table):
            table=shape.table
            tbl = table._tbl
            tobeRemoved=4-noOfComp
            reqHeight=0
            for row in table.rows:
                reqHeight+=row.height
            if(tobeRemoved!=0):
                for x in range(1,2*tobeRemoved+1):
                    tr = table.rows[11-x]._tr
                    tbl.remove(tr)
            for row in table.rows:
                row.height = reqHeight//(1+2*(1+noOfComp))
            
            for brand in sortedBrands:
                run = (table.cell(cnt,0)).text_frame.paragraphs[0].runs[0]
                run.text = run.text.replace(run.text,brand)
                cnt+=2
            for r in range(1,2*(noOfComp+1)+1):
                for c in range(1,3):
                    temp = arr[r-1][c-1]
                    if(c==2):
                        temp = str(temp)+"%"
                    run = (table.cell(r,c)).text_frame.paragraphs[0].runs[0]
                    run.text = run.text.replace(run.text,temp)

def getDesc(ids):
    print(ids)
    url = f'https://www.googleapis.com/youtube/v3/videos?part=snippet&id={ids}&key={apiKey}'
    response = requests.get(url)
    resObj = response.json()
    return resObj["items"][0]['snippet']['title']

def getViews(id):
    print(id)
    url = f'https://www.googleapis.com/youtube/v3/videos?part=statistics&id={id}&key={apiKey}'
    response = requests.get(url)
    resObj = response.json()
    return resObj["items"][0]['statistics']['viewCount']

def fetchData12(excelPathVideo,mainBrand,slideNo,sortedBrands):
    df = pd.read_excel(excelPathVideo,"video")
    pt = pd.pivot_table(df,values = ["VIEWS"],
                        index=["IAB_TIER2","VIDEO_ID"],
                        aggfunc={"VIEWS":"sum"},
                        columns=["CHANNEL_NAME"]  
                        )
    pt.columns=pt.columns.droplevel(level=0)
    subcats=[]
    arr=[]
    finaldf = pd.DataFrame()
    for brand in sortedBrands:
        subcat = pt[brand].dropna().groupby("IAB_TIER2").sum().sort_values(ascending=False).reset_index().iloc[0].iloc[0]
        subcatGrandTotal = pt[brand].dropna().groupby("IAB_TIER2").sum().sort_values(ascending=False).reset_index().iloc[0].iloc[1]
        subcats.append(subcat)
        tempdf = pt[brand].dropna().loc[subcat].sort_values(ascending=False).reset_index().loc[:1]
        tempdf['_Per']=((tempdf[brand]/subcatGrandTotal)*100).round(2)
        tempdf['_Desc']=[getDesc(tempdf['VIDEO_ID'].iloc[0]),getDesc(tempdf['VIDEO_ID'].iloc[1])]
        print(tempdf)
        finaldf=pd.concat([finaldf,tempdf],axis=0)
    arr = finaldf.loc[:,['_Per','_Desc']].to_numpy()
    slide = presentation.slides[slideNo-1]
    cnt=1
    for shape in slide.shapes:
        if(shape.has_table):
            table=shape.table
            tbl = table._tbl
            tobeRemoved=4-noOfComp
            reqHeight=0
            for row in table.rows:
                reqHeight+=row.height
            if(tobeRemoved!=0):
                for x in range(1,2*tobeRemoved+1):
                    tr = table.rows[11-x]._tr
                    tbl.remove(tr)
            for row in table.rows:
                row.height = reqHeight//(1+2*(1+noOfComp))
            for brand in sortedBrands:
                run = (table.cell(cnt,0)).text_frame.paragraphs[0].runs[0]
                run.text = run.text.replace(run.text,brand)
                cnt+=2
            cnt=1
            for subcat in subcats:
                run = (table.cell(cnt ,1)).text_frame.paragraphs[0].runs[0]
                run.text = run.text.replace(run.text,subcat)
                cnt+=2
            for r in range(1,2*(noOfComp+1)+1):
                for c in range(0,2):
                    temp = arr[r-1][c-1]
                    if(c==1):
                        temp = str(temp)+"%"
                    run = (table.cell(r,c+2)).text_frame.paragraphs[0].runs[0]
                    run.text = run.text.replace(run.text,temp)

def getIds(sortedBrands):
    df = pd.read_excel(excelPathVideo,"video")
    pt = pd.pivot_table(df,values = ["VIEWS"],
                        index=["VIDEO_ID"],
                        aggfunc={"VIEWS":"sum"},
                        columns=["CHANNEL_NAME"]  
                        )
    pt.columns=pt.columns.droplevel(level=0)
    ids=[]
    for brand in sortedBrands:
        x=1
        if(brand==mainBrand):
           x=2 
        df = pt.sort_values(brand,ascending=False)[brand].dropna().reset_index().iloc[:x].to_numpy()
        dflist = df.tolist()
        dflist[0].insert(2,brand)
        ids.append(dflist[0])
        if(brand==mainBrand):
            dflist[1].insert(2,brand)
            ids.append(dflist[1])
    return ids

def downloadImage(sortedBrands):
    idList = getIds(sortedBrands)
    for ids in idList:
        id=ids[0]
        link = f"https://i.ytimg.com/vi/{id}/maxresdefault.jpg"
        altLink = f"https://i.ytimg.com/vi/{id}/sddefault.jpg"
        outputPath=os.path.join("Images")
        fileName=id+".jpg"
        try:
            file_path = wget.download(link,out=outputPath)
            new_file_path = os.path.join(outputPath, fileName)
            os.rename(file_path, new_file_path)
        except:
            try:
                file_path = wget.download(altLink,out=outputPath)
                new_file_path = os.path.join(outputPath, fileName)
                os.rename(file_path, new_file_path)
            except:
                print("image not found")

def compact_number(number):
    suffixes = ['', 'K', 'M', 'B', 'T']
    if number < 1000:
        return str(number)
    exp = int((len(str(abs(number))) - 1) / 3)
    return '{:.1f}{}'.format(number / (1000 ** exp), suffixes[exp])

def fetThumbnails(sortedBrands):
    slide = presentation.slides[12]
    downloadImage(sortedBrands)
    idList = getIds(sortedBrands)
    cnt=0
    cnt2=0
    for shape in slide.shapes:
        if(hasattr(shape, "text") and cnt2<(2+noOfComp)):
            views = compact_number(int(getViews(idList[cnt2][0])))
            for para in shape.text_frame.paragraphs:
                if(para.runs):
                    for run in para.runs:
                        run.text = run.text.replace(run.text, views+" Views")
            cnt2+=1
    for shape in slide.shapes:
        if(shape.shape_type == 13 and cnt<(2+noOfComp)):
            img_shape=shape
            id=idList[cnt][0]
            fileName = id+".jpg"
            try:
                op = os.path.join("Images",fileName)
                slide.shapes.add_picture(op, img_shape.left, img_shape.top,img_shape.width,img_shape.height)
                slide.shapes._spTree.remove(img_shape._element)
            except:
                print("file not found")
            cnt=cnt+1
            
def fetchData89(excelPathChannel,mainBrand,slideNo,sortedBrands,type):
    df = pd.read_excel(excelPathChannel,"channel")
    pt = pd.pivot_table(df,values = [type],
                        index=["Date"],
                        aggfunc={type:"sum"},
                        columns=["CHANNEL_NAME"]  
                        )
    pt=pt.dropna()
    pt.columns=pt.columns.droplevel(level=0)
    pt=pt.diff().dropna()
    
    slide = presentation.slides[slideNo-1]
    for shape in slide.shapes:
        if(shape.has_chart):
            chart=shape.chart
            chart_data = ChartData()
            chart_data.categories = pt.index
            for brand in sortedBrands:
                chart_data.add_series(brand,pt[brand])
            chart.replace_data(chart_data)

def prepareDeck():

    #overall
    changeTextType("Century Gothic")

    #slide1
    addPptLogo(pptLogo,0,1) #location, slide-1, imageNumber-1

    #slide5
    sortedBrands = fetchData5(excelPathChannel,mainBrand,5) #excelPathChannel, mainbrandName, Slide
    addComa(4,1,noOfComp+1,1,5) #slideNumber-1, rowStart-1, rowEnd-1, colStart-1, colEnd-1

    fetchData6(excelPathChannel,mainBrand,6,sortedBrands)
    fetchData7(excelPathChannel,mainBrand,7,sortedBrands)
    fetchData89(excelPathChannel,mainBrand,8,sortedBrands,"TRUEVIEW_SPEND_ESTIMATE")
    fetchData89(excelPathChannel,mainBrand,9,sortedBrands,"VIEWS")
    fetchData10(excelPathChannel,mainBrand,10,sortedBrands)
    fetchData11(excelPathVideo,mainBrand,11,sortedBrands)
    fetchData12(excelPathVideo,mainBrand,12,sortedBrands)
    fetThumbnails(sortedBrands)



    # addComa(13,1,10,1,2) #slideNumber-1, rowStart-1, rowEnd-1, colStart-1, colEnd-1

    presentation.save(os.path.join("test_pixability.pptx"))

    send()
